import streamlit as st
from langgraph.graph import StateGraph, END
from langchain.chains import ConversationalRetrievalChain, LLMChain
from langchain_community.vectorstores import FAISS
from langchain_community.document_loaders import PyPDFLoader
from langchain.memory import ConversationBufferMemory
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate, PromptTemplate
from langchain_community.utilities import SQLDatabase
from typing import TypedDict, Optional
from pptx import Presentation
from pptx.util import Inches
import sys 
import sqlite3
import pandas as pd
import datetime

if 'initialized' not in st.session_state:
    st.session_state.initialized = False
    st.session_state.chat_history = []

# Predefined API key and vector store paths

groq_api_key = st.secrets["GROQ_API_KEY"]

embedding = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
api_key = groq_api_key
llm = ChatGroq(
    groq_api_key=api_key,
    model_name="llama-3.3-70b-versatile"
)

@st.cache_resource
def create_qa_chain(pdf_path):
    loader = PyPDFLoader(pdf_path)
    pages = loader.load_and_split()
    vectordb = FAISS.from_documents(pages, embedding)

    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectordb.as_retriever(search_kwargs={"k": 4}),
        memory=memory
    )
    return qa_chain

finance_prompt = PromptTemplate.from_template("""
You are a financial data extraction expert. Your job is to answer questions about financial data and ALWAYS present the answers in a clear markdown table format.

INSTRUCTIONS:
1. Extract relevant financial data from the provided context
2. Structure ALL responses as markdown tables with proper headers
3. Use numerical formatting for financial figures (e.g., currency, percentages)
4. Include column headers that clearly describe the data
5. Include a markdown table separator row with dashes after the header row

Context:
-----------------
{context}
-----------------

Chat History:
{chat_history}

Question: {question}

YOUR ANSWER MUST BE IN MARKDOWN TABLE FORMAT like this:
| Header1 | Header2 | Header3 |
| --- | --- | --- |
| Value1 | Value2 | Value3 |
| Value4 | Value5 | Value6 |

Answer:
""")
@st.cache_resource
def create_finance_chain():
    print("# Create the finance chain with custom prompt")
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)
    loader = PyPDFLoader("data/Finance.pdf")
    pages = loader.load_and_split()
    vectordb = FAISS.from_documents(pages, embedding)

    agent_d_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectordb.as_retriever(search_kwargs={"k": 4}),
        memory=memory,
        combine_docs_chain_kwargs={"prompt": finance_prompt}
    )
    return agent_d_chain

agent_d_chain = create_finance_chain()
@st.cache_resource
def create_chain():
    print("Chain is called")

    # Load FAISS vector store
    vectorstore = FAISS.load_local("faiss_index", embeddings=embedding,allow_dangerous_deserialization=True )
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    qa_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 4}),
        memory=memory
    )
    print("QA Chain", qa_chain)
    return qa_chain


# State schema
class AgentState(TypedDict):
    question: str
    answer: Optional[str]
    ppt: Optional[bool]
    ppt_file_content: Optional[str]
    filename: Optional[str]
    excel : Optional[bool]
    excel_path : Optional[str]
    excel_name : Optional[str]
    agent_type : Optional[str]


# Create agents with their own PDFs
agent_a_chain = create_qa_chain("data/Clinical_Trial.pdf")
agent_b_chain = create_qa_chain("data/Radiotherapy.pdf")  # Agent about Radiotherapy
agent_c_chain = create_qa_chain("data/IP_Manual.pdf")

def extract_slide_content(text):
    prompt = f"""
    Extract key points from the following document that would be useful in a presentation.
    Output should be in the form:
    - Title Slide: [Title]
    - Slide 1: [Title] | [Bullet Points] Always return the bullet points with * in the beginning.
    - Slide 2: [Title] | [Bullet Points] [TABLE:] If any tabular data is present return it in tabular format with TABLE: in the beginning.
    ...

    Remember : For cost related details return data always in tabular format.
    Always return the details in the above format only.
    Document:
    {text}
    """

    prompt = PromptTemplate.from_template(prompt)
    formatted_prompt = prompt.format(text=text)
    response = llm.invoke(formatted_prompt)
    selected = response.content.strip()

    return selected

def parse_slide_content(text):
    lines = text.strip().split('\n')
    slides = []
    current_slide = None
    in_table = False
    table_rows = []

    for line in lines:
        line = line.strip()
        if line.startswith("- Title Slide:"):
            current_slide = {
                "title": line.replace("- Title Slide:", "").strip(),
                "bullets": [],
                "table": [],
                "is_title": True
            }
            slides.append(current_slide)

        elif line.startswith("- Slide"):
            if current_slide and in_table:
                current_slide["table"] = table_rows
                table_rows = []
                in_table = False

            parts = line.split(":")
            title_and_content = parts[1].split("|")
            title = title_and_content[0].strip()
            current_slide = {
                "title": title,
                "bullets": [],
                "table": [],
                "is_title": False
            }
            slides.append(current_slide)
        elif line.startswith("•"):
            current_slide["bullets"].append(line.replace("•", "").strip())

        elif line.startswith("*"):
            current_slide["bullets"].append(line.replace("*", "").strip())

        elif line.startswith("TABLE:"):
            in_table = True
            table_rows = []

        elif in_table and "|" in line:
            row = [cell.strip() for cell in line.strip("|").split("|")]
            table_rows.append(row)

    if current_slide and in_table:
        current_slide["table"] = table_rows

    create_ppt_from_parsed_slides(slides=slides)

def create_ppt_from_parsed_slides(slides,template_path = "Theme.pptx",  output_file="slide_deck.pptx"):
    prs = Presentation(template_path)

    for slide in slides:
        layout = prs.slide_layouts[0] if slide["is_title"] else prs.slide_layouts[1]
        ppt_slide = prs.slides.add_slide(layout)
        ppt_slide.shapes.title.text = slide["title"]

        # Add bullets
        if not slide["is_title"] and slide["bullets"]:
            body = ppt_slide.placeholders[1]
            body.text = "\n".join(slide["bullets"])

        # Add table
        if slide["table"]:
            rows = len(slide["table"])
            cols = len(slide["table"][0])
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            height = Inches(0.8)

            table = ppt_slide.shapes.add_table(rows, cols, left, top, width, height).table

            for i, row in enumerate(slide["table"]):
                for j, cell in enumerate(row):
                    table.cell(i, j).text = cell

    prs.save(output_file)
    print(f"Slide deck saved to: {output_file}")

agent_e_chain = create_chain()
def extract_table_data(text):
    """
    Extract table data from markdown-formatted text and convert to pandas DataFrame.
    """
    import pandas as pd
    
    # Find table lines (lines starting with |)
    table_lines = [line.strip() for line in text.split('\n') if line.strip().startswith('|')]
    
    if len(table_lines) < 2:  # Need at least header and one data row
        return None
    
    # Extract headers
    header_line = table_lines[0]
    headers = [cell.strip() for cell in header_line.split('|') if cell.strip()]
    
    # Skip separator line if present (line with dashes like | --- | --- |)
    data_start_idx = 1
    if len(table_lines) > 1 and all(cell.strip().replace('-', '').replace(':', '') == '' 
                                   for cell in table_lines[1].split('|') if cell.strip()):
        data_start_idx = 2
    
    # Extract data rows
    data = []
    for line in table_lines[data_start_idx:]:
        # Split by | and remove empty cells at the beginning/end
        cells = [cell.strip() for cell in line.split('|')]
        cells = [cell for cell in cells if cell]  # Remove empty strings
        
        if cells:
            # Make sure we have the same number of cells as headers
            if len(cells) < len(headers):
                cells.extend([''] * (len(headers) - len(cells)))
            elif len(cells) > len(headers):
                cells = cells[:len(headers)]
                
            data.append(cells)
    
    # Create DataFrame
    if data:
        df = pd.DataFrame(data, columns=headers)
        
        return df
    
    return None


# Agent functions
def agent_a(state):
    # print("The ML agent is answering...")
    result = agent_a_chain({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "Clinical Trial Agent"}

def agent_b(state):
    # print("The Medical agent is answering...")
    result = agent_b_chain({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "Medical Agent"}

def agent_c(state):
    result = agent_c_chain({"question": state["question"]})
    return {"answer": result["answer"], "agent_type" : "IP Manual Agent"} 
def agent_d(state):
    """
    Finance agent that extracts data from finance PDF, formats it as a table,
    and saves it to Excel.
    """
    print("The Finance agent is answering...")
    
    # Get the question
    question = state["question"]
    
    # Modify question to explicitly request tabular data
    table_question = f"{question} (FORMAT RESPONSE AS A MARKDOWN TABLE WITH COLUMN HEADERS)"
    
    # Get answer from the chain
    result = agent_d_chain({"question": table_question})
    answer = result["answer"]
    
    # Check if the answer has a table format (contains | characters and header separator)
    if '|' in answer and any(line.strip().startswith('|') and line.strip().endswith('|') for line in answer.split('\n')):
        try:
            # Extract the table data to pandas DataFrame
            df = extract_table_data(answer)
            
            if df is not None and not df.empty:
                # Save the DataFrame to Excel
                excel_path = 'finance_data.xlsx'
                df.to_excel(excel_path, index=False)
                print(f"Successfully saved data to {excel_path}")
                
                # Add a note to the answer about the Excel file
                answer += f"\n\nThe data has been saved to {excel_path}"
        except Exception as e:
            print(f"Error processing table data: {e}")
    else:
        print("Response did not contain proper table format")
    
    return {"answer": answer, "excel": True, "excel_path" : "finance_data.xlsx", "excel_name" : "finance_data.xlsx", "agent_type" : "Finance Agent"}

def agent_e(state):
    result = agent_e_chain({"question": state["question"]})
    print("result",result)
    text = result["answer"]
    slide_content = extract_slide_content(text)
    print("Slide Content:", slide_content)
    # create_presentation(slide_content, filename="presentation.pptx")
    parse_slide_content(slide_content)
    response = {"answer": "The created slide deck is.", "ppt":True, "ppt_file_content":"slide_deck.pptx","filename":"presentation.pptx", "agent_type" : "Slide Deck Agent"}
    print("Resp",response)
    return response

def agent_f(state):
    db = SQLDatabase.from_uri("sqlite:///./Chinook.db")  # Replace with your DB URI
    llm = ChatGroq(
    groq_api_key=api_key,
    model_name="llama-3.3-70b-versatile",
    temperature = 0.4
    )
    prompt = ChatPromptTemplate.from_template(
        """
    You are a helpful assistant that translates natural language into SQL queries.
    Use the following database schema:

    {schema}

    Translate the following question into a SQL query:

    Question: {question}

    SQL Query:

    Return only the SQL query in response and nothing else. 
    Don't add any backticks in response
    """
    )
    chain = LLMChain(llm=llm, prompt=prompt)
    schema = db.get_table_info()
    sql_query = chain.run({"schema": schema, "question": state["question"]})
    print("Generated SQL Query:\n", sql_query)
    conn = sqlite3.connect('Chinook.db')
    df = pd.read_sql_query(sql_query, conn)

    return {"answer": df.to_markdown(index=False), "agent_type" : "SQL Agent"}

def error_handler_agent(state):
  print("Error handler agent is answering...")
  return {"answer": "The question seems to be out of context, please ask something else.", "agent_type" : "Error Handler Agent"}


#This can be cosidered as a knowledge base for the router agent.

agent_descriptions = {
    "agent_a": "Implement a Phase II, randomized, double-blind, placebo-controlled clinical trial to assess the efficacy and safety of XYZ Drug in patients with type 2 diabetes. The trial aims to evaluate the effect of XYZ Drug on HbA1c levels at 24 weeks.",
    "agent_b": "An expert in answering question related to Radiotherapy.",
    "agent_c": """Task Assignment: Investigational Product (IP) Management for XYZ Drug
                    As an agent, you are responsible for managing the Investigational Product (IP) for the XYZ Drug clinical trial. Your tasks include:
                    1. Product Handling
                    2. Dispensing
                    3. Record Keeping
                    4. Return and Disposal
                    5. Emergency Unblinding
                    Sponsor
                    Storage Conditions
                    Shelf Life""",
                    
    "agent_d": "A financial data expert that extracts and presents financial information in table format and saves data to Excel. Handles queries related to financial reports, metrics, and data extraction.",
    "agent_e": """An expert in creating powerpoint presentations or slide decks related to I2E Consulting.
     Fetch all the necessary details to include in the slide deck. Try top explain each slide with multiple bulllet points to make the slide deck more interactive.""",
    "agent_f" : "A sql expert agent which is capable of answering questions related to Album, Artist, Customer, Employee, Genre, Invoice, InvoiceLine, MediaType, Playlist, PlaylistTrack, Track as these are the tables present in the database."

}

from langchain.prompts import PromptTemplate

def supervisor_router(state):
    question = state["question"]

    system_prompt = """
You are a router for a multi-agent system. Given a user question and a list of available agents with their descriptions,
choose the best agent to handle the query.


If any of the user question is not directly related to any of the agents, return "No agent selected".

For response only return the Agent name and nothing else.

Agents:
{agent_list}

User question:
{question}
"""

    prompt = PromptTemplate.from_template(system_prompt)
    formatted_prompt = prompt.format(
        agent_list="\n".join([f"{name}: {desc}" for name, desc in agent_descriptions.items()]),
        question=question
    )

    response = llm.invoke(formatted_prompt)
    selected = response.content.strip().lower()
    print("Understanding of the router agent:", selected)

    # Sanitize response to match node names
    if "agent_a" in selected:
        return "agent_a"
    elif "agent_b" in selected:
        return "agent_b"
    elif "agent_c" in selected:
        return "agent_c"
    elif "agent_d" in selected:
        return "agent_d"
    elif "agent_e" in selected:
        return "agent_e"
    elif "agent_f" in selected:
        return "agent_f"
    else:
        print("Not a correct question")
        return "error_handler_agent"
@st.cache_resource
def chatbot_func():
    print("Running Chatbot func:")
    
    #LanggrapH Workflow
    from langgraph.graph import StateGraph, END

    workflow = StateGraph(state_schema=AgentState)
    workflow.add_node("supervisor", supervisor_router)
    workflow.add_node("agent_a", agent_a)
    workflow.add_node("agent_b", agent_b)
    workflow.add_node("agent_c", agent_c)
    workflow.add_node("agent_d", agent_d)
    workflow.add_node("agent_e", agent_e)
    workflow.add_node("agent_f", agent_f)
    workflow.add_node("error_handler_agent", error_handler_agent)
    workflow.set_conditional_entry_point(supervisor_router)


    workflow.add_edge("agent_a", END)
    workflow.add_edge("agent_b", END)
    workflow.add_edge("agent_c", END)
    workflow.add_edge("agent_d", END)
    workflow.add_edge("agent_e", END)
    workflow.add_edge("agent_f", END)
    workflow.add_edge("error_handler_agent", END)
    return workflow.compile()
 

graph = chatbot_func()
# Streamlit UI
st.title("Multi-Bot System")
st.write("Ask a question and get answers from specialized agents.")
query = st.text_input("Enter your question:")

if st.button("Get Answer"):
    if not query.strip():
        st.warning("Please enter a question before clicking the button.")
    else:
        print("query is:", query)
        result = graph.invoke({"question": query})

        agent_type = result.get("agent_type", None)
        if agent_type:
            st.markdown(
    f"<h3 style='color: #2b01be;'><small>Answer from the </small><b style='font-weight:900;'>{agent_type}</b></h3>",
    #f"<h3 style='color: #3A7CA5;'>Answer from the <b style='font-weight: 900;'>{agent_type}</b></h3>",
    # f"<h3 style='color: #3A7CA5;'>Answer from the <b>{agent_type}</b></h3>",
    unsafe_allow_html=True
)

            # st.subheader(agent_type)

        # 2) Always show the textual answer
        answer = result.get("answer", "No answer returned.")
        st.write(answer)

        # 3) If your agents have generated a PPT, show a download button
        if result.get("ppt", False):
            ppt_path = result.get("ppt_file_content")
            ppt_filename = result.get("filename", "presentation.pptx")
            try:
                with open(ppt_path, "rb") as f:
                    ppt_bytes = f.read()
                st.download_button(
                    label="Download PPT",
                    data=ppt_bytes,
                    file_name=ppt_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"Could not open PPT at {ppt_path}: {e}")

        # 4) If your agents have generated an Excel file, show a download button
        if result.get("excel", False):
            excel_path = result.get("excel_path")
            excel_filename = result.get("excel_name", "data.xlsx")
            try:
                with open(excel_path, "rb") as f:
                    excel_bytes = f.read()
                st.download_button(
                    label="Download Excel",
                    data=excel_bytes,
                    file_name=excel_filename,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
            except Exception as e:
                st.error(f"Could not open Excel at {excel_path}: {e}") 
